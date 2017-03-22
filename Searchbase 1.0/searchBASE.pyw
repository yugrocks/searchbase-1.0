from lxml import _elementpath as _dummy
from pptx import Presentation
import openpyxl
import pickle
import docx
from functools import*
import os
from tkinter import*
import collections.abc
import threading
lkill=False
frame1=None
label2=None
myButton=None
x={}
g=None
b=None
root1=None
root=None
box2=None
box4=None
box6=None
grint=False
gtxt=None
frameexists=False
upbox=False
def capspermutation(string):

   c=[]
   l=[]
   st3=string.upper()
   st2=string.lower()
   c.append(st2)
   c.append(st3)
   for alphabet in st2:
       l.append(alphabet)
   wrd=l[0].upper()
   e=1

   while e<len(l):
      wrd=wrd+l[e]
      e+=1
   c.append(wrd)
   
   return(c)

extlist=('.mp3','.ace','.mp4','.avi','.flv','.mkv','.jpg','.jpeg','.htm','.html','.pdf','.png','.exe','.txt','.docx','.xml','.gif','.odp','.pptx','.zip','.py','.cpp','.c','.rar','.xlsx')


def database(dir):
    
    b=os.listdir(dir)
    for entity in b:
        if entity=="$Recycle.Bin" or entity=="$RECYCLE BIN" or entity=="System Volume Information":
           continue
        if os.path.isdir(r'{}\{}'.format(dir,entity)):
           x[entity]=r'{}\{}'.format(dir,entity)
           try:
             database(r'{}\{}'.format(dir,entity))
           except:
             continue
        else:
           if '.mp3' in entity:
              if '.lnk' in entity or '.LNK' in entity:
                 continue                
              meta=getmeta(r'{}\{}'.format(dir,entity))
              tuple2=(entity,meta)
              x[tuple2]=r'{}\{}'.format(dir,entity)
           else:
              if '.lnk' in entity or '.LNK' in entity:
                 continue   
              for ext in extlist:
                 if ext in entity:
                    x[entity]=r'{}\{}'.format(dir,entity)

def getmeta(file):
   os.chdir(r'C:\database1')
   try:
      a=os.popen(r'tool.bat "{}"'.format(file))
   except:
      createtool()
      a=os.popen(r'tool.bat "{}"'.format(file))
   for  word in a:
      if 'artist' in word or 'Artist' in word:
         return word


updating=False
thislabel=None
from tool import*               
updatecount=0
def update_database():
   global box2,upbox,box4,updatecount,frame2,framen,button5,binitial,fame1
   global myButton,root,entryBox,thislabel,myButtontwo,updating
   updating=True
   if button5!=None:
       button5.destroy()
   if binitial!=None:
       binitial.destroy()
   entryBox.destroy()
   myButton.destroy()
   myButtontwo.destroy()

   try:
      dingdong=Label(frame2,text='SCANNING...',bg='black',fg='blue')
      dingdong.pack(side=TOP)
   except:
      dingdong=Label(framen,text='SCANNING...',bg='black',fg='blue')
      dingdong.pack(side=TOP)
   if not os.path.exists(r'C:\database1'):
      os.mkdir(r'C:\database1')
   createtool()
   for direc in ('D:','E:','F:','G:','H:','I:','B:','A:'):
      try:
         rcheck=os.popen("vol {}".format(direc))
         rchk=rcheck.readline()
         if 'recovery' in rchk or 'RECOVERY' in rchk:
            continue
         
         database(direc)
      except:
         continue
   os.chdir(r'C:\Users')
   
   dirinfo=os.popen('dir /B')
   for direc in dirinfo:
      str(direc)
      dirc=direc.strip()
      try:
         database(r"C:\Users\{}".format(dirc))
      except:
         continue

   update_file_database()
   if frame2!=None:
      frame2.destroy()
   if upbox:
      framen.destroy()
   dingdong.destroy()
   createTextBox(root)
   myButton = Button(root, text="Search",bg="SteelBlue1",relief=FLAT,activebackground="gray53",width=14)
   myButton.pack(side=TOP,pady=1)
   Label(root,text="We Are Good To Go!",fg='white',bg='black').pack(side=TOP)
   thislabel=Label(root,text="The Documents are being scanned and indexed. May take a while. Plz dont exit during this operation.",fg='white',bg='blue')
   thislabel.pack()
   thfun2.start()


thfun=threading.Thread(target=update_database)

def update_file_database():
   global updatecount
   if not os.path.exists(r'C:\database1'):
       os.mkdir(r'C:\database1')
   os.chdir(r'C:\database1')
   fileopen=open('DATABASEF.pkl','wb')
   pickle.dump(x,fileopen)
   fileopen.close()


def opennow(path):
    os.popen('explorer.exe "{}"'.format(path))
    root.destroy()
    os._exit(0)


frame2=None
button5=None
def upcommand():
    global frame2,root,button5,frame1,updating,frameexists
    updating=True
    frame2=Frame(root)
    frame2.pack(side=TOP)
    frame2.config(background='black')
    frame1.destroy()
    frameexists=False
    Label(frame2,text="   - Updating, A Matter Of Seconds..",bg='black',fg='white').pack(side=TOP)
    Label(frame2,text="   - Everytime it is updated, old records are deleted.",bg='black',fg='white').pack(side=TOP)
    Label(frame2,text=r"   - In drive C:, the custom folder C:\Users\whtever  will be scanned and your custom drive(if any) for ex. F: or whatever it is named will be scanned",bg='black',fg='white').pack(side=TOP)
    Label(frame2,text="   - Plz Be Patient! I will inform you when done.",bg='black',fg='white').pack(side=TOP)
    button5=Button(frame2,text="Do it!",command=thfun.start,bg='dim gray')
    button5.pack(ipadx=20)



def remove(s):
   b=[',','<','.','>','/','\',','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(','\n',')','+','=','...',"{","}","|","\\"]
   c=s.strip()
   n=len(c)
   list1=[]
   list2=[]
   for i in range(0,n):
       list1.append(c[i])
   d=0
   for i in range(0,n):
       if list1[d] in b:
          list2.append(list1[d])
          list1[d]=' '
       d=d+1
   return list1,c

not_destroyed=True

def splt(string):
    e=''
    s=remove(string)
    d=[s[1]]
    for alpha in s[0]:
        e+=alpha
    b=e.split()
    for element in b:
        d.append(element)
    listcaps=capspermutation(string)
    if len(b)==1:
       listcaps1=capspermutation(b[0])
       for element in listcaps1:
         d.append(element)
    if len(b)==2:
       listcaps2=capspermutation(b[1])
       for element in listcaps2:
         d.append(element)  
    if len(b)==3:
       listcaps3=capspermutation(b[2])
       for element in listcaps3:
         d.append(element)
    for element in listcaps:
        d.append(element)
    return d

framen=None
u=0
binitial=None
def search(key,frame1):
    global box2,framen,binitial
    global u
    y={}
    global frameexists
    global upbox
    frameexists=True
    if searchlabel!=None:
        searchlabel.destroy()
    try:
       os.chdir(r'C:\database1')
       fileopen=open('DATABASEF.pkl','rb')
       dbse=pickle.load(fileopen)
       checktest=True

    except:
       upbox=True
       frame1.destroy()
       frameexists=False
       framen=Frame(root,bg='black')
       framen.pack(side=TOP)
       label=Label(framen,text='Please Update The Database To get started. This may take a while. Please wait until the application starts to respond properly',bg='black',fg='white')
       label.pack(side=TOP)
       Label(framen,text='   - A Matter Of Seconds, A One Time Investment.',bg='black',fg='white').pack(side=TOP)
       Label(framen,text="   - Everytime it is updated, old records are deleted.",bg='black',fg='white').pack(side=TOP)
       Label(framen,text=r"  - In drive C:, the custom folder C:\Users\whatever  will be scanned and your custom drive(if any) for ex. F: or whatever it is named will be scanned",bg='black',fg='white').pack(side=TOP)
       Label(framen,text="   - Plz Be Patient. It will show up when done.",bg='black',fg='white').pack(side=TOP)
    
       binitial=Button(framen,text='Update',command=thfun.start,bg='gray',activebackground='indian red',relief=FLAT)
       binitial.pack(side=TOP)
       
    list=splt(str(key))
    n=len(list)
    if checktest==True:
       for word in dbse:
           for alpha in list:
               if type(word) is str:
            
                  if alpha in word.lower():
                    if word in y:
                       continue
                    y[word]=dbse[word]
                    u+=1
               elif type(word) is tuple:
                  if alpha in word[0].lower():
                     if word in y:
                        continue
                     y[word]=dbse[word]
                     u+=1
                  if word[1]!=None:
                     if word in y:
                        continue
                     if alpha in word[1]:
                        y[word]=dbse[word]
                        u+=1
               if u==800:
                    break
       e=n
       k=0


       yo='No'
       for rsnt in y:
         if type(rsnt)is str:
            if key.lower() in rsnt.lower():
               if y[rsnt]=='':
                  continue
               yo='yes'
               dirnam=os.path.dirname(y[rsnt])
               button=Button(frame1,text=rsnt,command=partial(opennow,y[rsnt]),bg="SteelBlue2",relief=GROOVE,activebackground="gray53",height=2)
               button.pack(side=TOP)
               button2=Button(frame1,text="Open file location",command=partial(opennow,dirnam),bg="gray53",relief=FLAT,activebackground="SteelBlue2",height=1)
               button2.pack(side=TOP)
            
               k+=1
               if k==10:
                  break
            
               y[rsnt]=''

         if type(rsnt) is tuple:
         
            if key.lower() in rsnt[0].lower() :
               if y[rsnt]=='':
                   continue
               yo='yes'
               dirnam=os.path.dirname(y[rsnt])
               button=Button(frame1,text=rsnt[0],command=partial(opennow,y[rsnt]),bg="SteelBlue2",relief=GROOVE,activebackground="gray53",height=2)
               button.pack(side=TOP)
               button2=Button(frame1,text="Open file location",command=partial(opennow,dirnam),bg="gray53",relief=FLAT,activebackground="SteelBlue2",height=1)
               button2.pack(side=TOP)
            
               k+=1
               if k==10:
                  break
               y[rsnt]=''
            if rsnt[1]!=None:
               if key.lower() in rsnt[1].lower():
                  if y[rsnt]=='':
                      continue
                  yo='yes'
                  dirnam=os.path.dirname(y[rsnt])
                  button=Button(frame1,text=rsnt[0],command=partial(opennow,y[rsnt]),bg="SteelBlue2",relief=GROOVE,activebackground="gray53",height=2)
                  button.pack(side=TOP)
                  button2=Button(frame1,text="Open file location",command=partial(opennow,dirnam),bg="gray53",relief=FLAT,activebackground="SteelBlue2",height=1)
                  button2.pack(side=TOP)
            
                  k+=1
                  if k==10:
                     break

                  y[rsnt]=''
            
       for i in range(0,n):
           if yo=='yes':
              break
           for rsnt in y:

                   d=0
                   for z in list:
                      if type(rsnt) is tuple:
                         if z in rsnt[0].lower():
                            d+=1
                         if rsnt[1]!=None:
                            if z in rsnt[1].lower():
                               d+=1
                      if type(rsnt) is str:
                         if z in rsnt.lower():
                            d+=1


                   if d==e:
                       if y[rsnt]=='':
                           continue
                    
                       dirnam=os.path.dirname(y[rsnt])
                       if type(rsnt) is tuple:
                          button=Button(frame1,text=rsnt[0],command=partial(opennow,y[rsnt]),bg="gray53",relief=GROOVE,activebackground="SteelBlue2",height=2)
                          button.pack(side=TOP)
                          button2=Button(frame1,text="Open file location",command=partial(opennow,dirnam),bg="SteelBlue2",relief=FLAT,activebackground="gray53",height=1)
                          button2.pack(side=TOP)
                    
                       if type(rsnt) is str:
                          button=Button(frame1,text=rsnt,command=partial(opennow,y[rsnt]),bg="gray53",relief=GROOVE,activebackground="SteelBlue2",height=2)
                          button.pack(side=TOP)
                          button2=Button(frame1,text="Open file location",command=partial(opennow,dirnam),bg="SteelBlue2",relief=FLAT,activebackground="gray53",height=1)
                          button2.pack(side=TOP)
                       y[rsnt]=''
                       k=k+1
                       if k==8:
                          break
           e-=1
       global lkill,label2
       if k==0:
           label2=Label(root,text="No results",bg="SteelBlue1",height=3,width=30)
           label2.pack()
           label2.place(height=20,width=90)
           lkill=True
       if k!=0 and lkill==True:
           label2.destroy()
           yo='No'

def clearf(frame):
    a=frame.winfo_children()
    for widget in a:
        widget.destroy()
txt=None
itxt=None
def buttonPushed():
   global entryBox,txt,itxt,updating
   global root
   global frame1
   global frameexists
   if txt!=None:
      itxt=txt
   txt = entryBox.get()
   if txt!='' and itxt!=txt and txt!=' ' and not updating:
        if frameexists:
           info=frame1.winfo_children()
           for widget in info:
              widget.destroy()
        search(txt,frame1)

def buttonpushed2(event):
     global entryBox,txt,itxt,thinitbut
     global root
     global frame1
     global frameexists
     if txt!=None:
        itxt=txt
     txt = entryBox.get()
     if txt!='' and itxt!=txt and txt!=' ':
          if frameexists:
             info=frame1.winfo_children()
             for widget in info:
                widget.destroy()
          search(txt,frame1)



def initbut():
        global txt,mailinitiated,updating
        while True:
           try:
              if txt!=None and ' ' in txt or mailinitiated or updating:
                  labelenter=Label(root,text="Press Enter To Search",bg="SteelBlue1",height=3,width=30)
                  labelenter.pack()
                  labelenter.place(height=20,width=120)
                  break
              buttonPushed()
           except:
              continue





def createTextBox(parent):
   global entryBox
   entryBox = Entry(parent,background="lemon chiffon",width=60,relief=SUNKEN,font='Mincho')
   entryBox.pack(ipady=7,side=TOP)
   entryBox.focus_set()
   entryBox.bind('<Return>',buttonpushed2)


u=0
v=0
x={}
def checksum(dir):   #a function that returns the checksum of a given file of any format
    a=os.popen('certutil -hashfile "{}"'.format(dir))
    for word in a:
       b=a.readline()
       return b

completed=False
def check_dir(dire):
    global v,root1
    global u,completed

    b=os.listdir(os.chdir(dire))   #b is a list
    for entity in b:   #entity is the current folder or file being processed
        if os.path.isdir("{}\{}".format(dire,entity)):
            if len(os.listdir("{}\{}".format(dire,entity)))==0:
                os.rmdir("{}\{}".format(dire,entity))
                v+=1
                Label(root1,text="The folder {}\{} has been deleted because it was empty".format(dire,entity)).pack(side=TOP)
                continue
            check_dir("{}\{}".format(dire,entity))


        else:
            c=checksum("{}\{}".format(dire,entity))
            if c in x:
                x[c]+=1
                os.remove("{}\{}".format(dire,entity))
                u+=1
                #print("the file {}\{} has been deleted".format(dir,entity))
            else:
                x[c]=1
    completed=True
l1=None
l2=None
proceed=None
dupinitiated=False
def init_dup():
   global root1,b,proceed,dupinitiated
   root1=Tk()
   root1.title("Remove duplicate files")
   dupinitiated=True
   global box6


   Message(root1,text="WARNING: Please note that any empty folders and empty files will also be deleted. \nThe app may hand during this operation",bg='antique white',width=500).pack(side=TOP)
   proceed=Label(root1,text="Do you want to proceed? [Y/N]")
   proceed.pack(side=TOP)
   box6=Entry(root1)
   box6.pack(side=TOP)
   box6.focus_set()
   b=Button(root1,text='Go On',bg='gray53',activebackground='peach puff')
   b.pack()
   b.bind('<Button-1>',get_response)
   root1.mainloop()
but=None


def get_response2():
    global g,box6,root1,not_destroyed,but,l1,l2,root

    h=box6.get()
    x={}

    try:
        check_dir(h)   #calling the function
        box6.destroy()
        but.destroy()
        l1.destroy()
        l2.destroy()
        if not_destroyed:
           if u!=0 or v!=0:  #checking if any file was actually deleted
              Message(root1,text='Total {} duplicate file(s) and {} empty folder(s) was/were deleted.\nPlease Exit the whole application.\nPlease do not use this feature more than once in one session.'.format(u,v),width=300).pack(side=TOP)
           elif completed:
              Message(root1,text="Scan complete.\nNo duplicate files were found.\nNow please Exit the whole application.\nPlease do not use this feature more than once in one session.",width=300).pack(side=TOP)
        exitbu=Button(root1,text='EXIT',bg='gray53',activebackground="peachpuff")
        exitbu.pack(side=TOP)
        exitbu.focus_set()
        exitbu.bind('<Return>',dest)
        exitbu.bind('<Button-1>',dest)

    except:
        Label(root1,text="Path Not found.Operation failed.Please restart the application").pack(side=TOP)

def dest(event):
    os._exit(0)




def get_response(event):

    global box6,g,b,root1,u,v,but,proceed,grint,gtxt,l1,l2
    global not_destroyed
    grint=True
    box6.focus_set()
    g=box6.get()
    box6.delete(0,END)
    proceed.destroy()
    if g=='Y' or g=='y':
       l1=Label(root1,text="Enter the drive or directory(full path) in which")
       l2=Label(root1,text="you want to remove duplicated files.For Ex. 'F:'")
       l1.pack()
       l2.pack()

       b.destroy()
       b=None

       but=Button(root1,text='DO IT',bg='gray53',activebackground="peach puff",command=get_response2)
       but.pack(side=TOP)


    elif g=='N' or g=='n':
        root1.destroy()
    else:
        Label(root1,text="Unrecognized response. Please try again")
        root1.destroy()
        init_dup()
        not_destroyed=False


def removeit(s):
   b=[',','<','.','>','/','\',','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(','\n',')','+','=','...',"{","}","|","\\"]
   n=len(s)
   list1=[]
   for i in range(0,n):
       list1.append(s[i])
   c=0
   for i in range(0,n):
       if list1[c] in b:
          del list1[c]
          continue
       c=c+1
   return list1

yugbox=None
yugbox2=None
yugframe=None
def yughere(event):
   global yugbox,yugbox2,yugframe
   iamtrue=True
   yug1=str(yugbox.get())
   yug2=str(yugbox2.get())
   try:
      a=open(yug1)
      d=open(yug2,'x')
   except:
      Label(yugframe,text='Error: Either the path is incorrect or the second file already exists.').pack(side=TOP)
      iamtrue=False
   if iamtrue:
     try:
         for word in a:
          c=removeit(word)
          for l in c:
              d.write(l)

          d.write("\n")
     except:
         Label(yugframe,text='Something went wrong.Please Try again')
   if iamtrue:
      clearf(yugframe)
      Label(yugframe,text="Ok, I'm done! Let's now meet the unhappy file (:P) ",height=5,bg="SteelBlue2").pack()



def inityughere():
    global yugbox,yugbox2,yugframe
    root2=Tk()
    root2.config(background='RoyalBlue4')
    yugframe=Frame(root2,bg='RoyalBlue4')
    yugframe.pack()
    Label(yugframe,text="Enter the full path to the text file(.txt extension)",bg='RoyalBlue4',fg='white').pack(side=TOP)
    yugbox=Entry(yugframe,bg='lemon chiffon')
    yugbox.pack(side=TOP)
    Label(yugframe,text="Now enter the full path (including file name which will be created automatically) to store the 'unhappy' text",fg='white',bg='RoyalBlue4').pack(side=TOP)
    yugbox.focus_set()
    yugbox2=Entry(yugframe,bg='lemon chiffon')
    yugbox2.pack(side=TOP)
    yugbox.bind('<Return>',focusSET)
    yugbox2.bind('<Return>',yughere)
    yugbutton=Button(yugframe,text="Do This Crap!",bg='gray53')
    yugbutton.pack(side=TOP)
    yugbutton.bind('<Button-1>',yughere)

def focusSET(event):
   global yugbox2
   yugbox2.focus_set()


def removenew(s):
   thisb=[',','<','.','>','/','?',"'",'"',':',';','[',']','_','-','*','&','^','%','$','#','@','!','~','(',')','+','=','...',"{","}","|","\\"]
   n=len(s)
   thislist1=[]
   for i in range(0,n):
       thislist1.append(s[i])
   c=0
   for i in range(0,n):
       if thislist1[c] in thisb:
          del thislist1[c]

          continue
       c=c+1

   string=''

   for ele in thislist1:
      string+=ele
   thislist2=string.split()
   return thislist2


def scan2():
   global thislabel,updating
   updating=True
   stopwordlist=['None','is','an','the','are','a','of','and','to','for','in','it','',' ']

   thisfile1=open(r'C:\database1\DATABASEF.pkl','rb')
   dbse=pickle.load(thisfile1)
   xdct={}

   for entity in dbse:
     if '.docx' in entity:
       
        try:
           a=docx.Document(dbse[entity])
           for line in a.paragraphs:
              b=removenew(line.text)
              for word in b:
                 if word=='' or word==' ':
                    continue
                 elif '\u2019' in word:
                    wordn=word.replace('\u2019','')
                 elif '\u2026' in word:
                    wordn=word.replace('\u2026','')
                 else:
                    wordn=word
                 if wordn in xdct:
                    if dbse[entity] in xdct[wordn]:
                       continue
                    xdct[wordn].append(dbse[entity])
                    continue
                 listm=[]
                 listm.append(dbse[entity])
                 xdct[wordn]=listm

        except:
           continue

     if '.txt' in entity or '.cpp' in entity or '.csv' in entity:
        try:
           am=open(dbse[entity])

           for word in am:
               list2=removenew(word)
               for e in list2:
                  d=e.replace('\x00','')
                  if d=='' or d==' ' or '1' in d or '2' in d or '3' in d or '4' in d or '5' in d or '6' in d or '7' in d or '8' in d or '9' in d or '0' in d or d in stopwordlist:
                     continue
                  if d in xdct:
                     if dbse[entity] in xdct[d]:
                        continue
                     xdct[d].append(dbse[entity])
                     continue
                  listn=[]
                  listn.append(dbse[entity])
                  xdct[d]=listn

        except:
           continue

     if '.pptx' in entity:

        try:
            prs = Presentation(dbse[entity])
            for slide in prs.slides:
               for shape in slide.shapes:
                  if not shape.has_text_frame:
                     continue
                  for paragraph in shape.text_frame.paragraphs:
                     a=removenew(paragraph.text)
                     for ptword in a:
                        if ptword in stopwordlist:
                           continue
                        if ptword in xdct:
                            if dbse[entity] in xdct[ptword]:
                               continue
                            xdct[ptword].append(dbse[entity])
                            continue
                        listn=[]
                        listn.append(dbse[entity])
                        xdct[ptword]=listn
        except:
           continue

     if '.xlsx' in entity:
   
        try:
           wb=openpyxl.load_workbook(dbse[entity])
           print(entity)
           getnm=wb.get_sheet_names()
           for sheetn in getnm:
              sheet=wb.get_sheet_by_name(sheetn)
              print(sheet)
              hr=sheet.get_highest_row()
              hc=sheet.get_highest_column()
              
              for i in range(1,hr):
                 for j in range(1,hc):
                    clvalue=sheet.cell(row=i,column=j).value
                    print(clvalue)
                    if type(clvalue) is not int:
                       if clvalue in stopwordlist or clvalue is None or '1' in clvalue or '2' in clvalue or '3' in clvalue or '4' in clvalue or '5' in clvalue or '6' in clvalue or '7' in clvalue or '8' in clvalue or '9' in clvalue or '0' in clvalue:
                           continue
                    else:
                       continue
                       print("LOL")
                    if type(clvalue) is str:
                       clvalue2=removenew(clvalue)
                       for ant in clvalue2:
                           print
                           (ant)
                           if ant in xdct:
                               xdct[ant].append(dbse[entity])
                               continue
                           listn=[]
                           print(listn)
                           listn.append(dbse[entity])
                           xdct[ant]=listn
                           print(xdct[ant])
                        
        except:
            continue


       

        
   print('Done')
   thisfile2=open('C:\database1\DATABASEW.pkl','wb')
   pickle.dump(xdct,thisfile2)
   thisfile2.close()
   thisfile1.close()
   thislabel2=Label(text="Great! Now search inside of all kinds of text docs as well",fg='white',bg='blue')
   thislabel2.pack()
   thislabel.destroy()



import operator

thfun2=threading.Thread(target=scan2)

altsearch=False
suplist=None
def search2(event):
    global frame1,frameexists,thisdict2,altsearch,suplist
    dict3={}
    altsearch=True
    if frameexists:
       clearf(frame1)
    txt = entryBox.get()
    thislist=[]
    if txt !='':
       for element in txt.split():
           elelist=capspermutation(element)
           for _ in elelist:
               thislist.append(_)

       for words in thislist:
           try:
              if words in thisdict2 and words!='':
                  if words in dict3:
                      continue
                  dict3[words]=thisdict2[words]
           except:
               continue

       n=len(thislist)
       e=n
       suplist={}
       count=0
       for word in dict3:
           if len(dict3[word])<100:
              for ent in dict3[word]:
                  if ent in suplist:
                      suplist[ent]+=1
                      continue
                  suplist[ent]=1
       s2button1=Button(frame1,text='See All',command=seeall,bg='white',fg='gray',relief=GROOVE)
       s2button1.pack(side=TOP,ipadx=15)
       frame1.config(background='white')
       for key,value in sorted(suplist.items(),key=operator.itemgetter(1),reverse=True):
           Button(frame1,text=os.path.basename(key),command=partial(opennow,key),bg='blue',fg='white',width=50,activebackground='gray53',relief=GROOVE,font='Mincho').pack(side=TOP,ipady=6)
           count+=1
           if count==15:
               count=0
               break


def seeall():
    global suplist
    if len(suplist)!=0:
       root5=Tk()
       root5.title('See All')
       root5.config(background='peach puff')
    count=0
    for key,value in sorted(suplist.items(),key=operator.itemgetter(1),reverse=True):
           Button(root5,text=os.path.basename(key),command=partial(opennow,key),bg='blue',fg='white',width=40,activebackground='gray53',relief=GROOVE,font='Mincho 10').pack(side=TOP,ipadx=50)
           count+=1
           if count==29:
               count=0
               break




thisvar=None
thisdict2=None
def initiates2():
    global thisdict2,thisvar,frameexists
    try:
       os.chdir(r"C:\database1")
       thisvar=open(r"C:\database1\DATABASEW.pkl",'rb')
       thisdict2=pickle.load(thisvar)
    except:
       if frameexists:
          upcommand()
s2t=threading.Thread(target=initiates2)
s2t.start()





import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email.utils import COMMASPACE, formatdate
from email import encoders

def send_mail( send_from, send_to, subject, text, files=[], server="smtp.gmail.com", port=587, username='', password='', isTls=True):
    msg = MIMEMultipart()
    msg['From'] = str(send_from)
    msg['To'] = COMMASPACE.join(str(send_to))
    msg['Date'] = formatdate(localtime = True)
    msg['Subject'] = str(subject)

    msg.attach( MIMEText(text) )

    for f in files:
        part = MIMEBase('application', "octet-stream")
        part.set_payload( open(f,"rb").read() )
        encoders.encode_base64(part)
        part.add_header('Content-Disposition', 'attachment; filename="{0}"'.format(os.path.basename(f)))
        msg.attach(part)

    smtp = smtplib.SMTP(server, port)
    if isTls: smtp.starttls()
    smtp.login(username,password)
    smtp.sendmail(send_from, send_to, msg.as_string())
    smtp.quit()






root4=None
mframe=None
mailinitiated=False
var=None
def init_mail():
    global root4,mframe,mailinitiated,var
    mailinitiated=True
    root4=Tk()
    root4.title('Quick Email Facility')
    root4.config(background='SteelBlue1')
    mframe=Frame(root4)
    mframe.pack()
    mframe.config(bg='SteelBlue2')
    var=0
    Message(mframe,text="To use this quick email facility, please make sure that you have \n enabled the 'Allow less secure Apps' option in gmail account settings. \n\nMy account-> Sign in And security-> Allow less Secure apps \n\nYour mailing address and password will never be stored.\nIt is totally safe.",bg='SteelBlue2',fg='black',font='Times',width=600).pack()
    chk=Checkbutton(mframe,text='Email With Attachment',command=variable,bg='SteelBlue2')
    chk.pack()
    Button(mframe,text='Done !',command=go_mail,bg='gray53',activebackground='peach puff',relief=GROOVE).pack(ipadx=20)
    root4.mainloop()

def variable():
    global var
    if var==0:
        var=1

    else:
        var=0
ethread=None
e1=None;e2=None;e3=None;e4=None;e5=None;e6=None;eframe=None
def go_mail():
    global mframe,root4,e1,e2,e3,e4,e5,e6,eframe,var,ethread
    mframe.destroy()
    eframe=Frame(root4)
    eframe.pack()
    eframe.config(background='SteelBlue1')
    e1=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Email:')
    l1.pack(ipadx=60,ipady=4)
    e1.pack(ipadx=60,ipady=4)
    e2=Entry(eframe,background='peach puff',font='Times',show='*')
    l1=Label(eframe,text='Password')
    l1.pack(ipadx=60,ipady=4)
    e2.pack(ipadx=60,ipady=4)
    e3=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Reciever(s):')
    l1.pack(ipadx=60,ipady=4)
    e3.pack(ipadx=60,ipady=4)
    e4=Entry(eframe,background='peach puff',font='Times')
    l1=Label(eframe,text='Subject:')
    l1.pack(ipadx=60,ipady=4)
    e4.pack(ipadx=60,ipady=4)
    e5=Text(eframe,background='peach puff',font='Times',height=10,width=40)
    l1=Label(eframe,text='Body:')
    l1.pack(ipadx=60,ipady=4)
    e5.pack(ipadx=60,ipady=4)

    if var:
       e6=Entry(eframe,background='peach puff',font='Times')
       l1=Label(eframe,text='Attachment file(Full Path To The File)')
       l1.pack(ipadx=60,ipady=4)
       e6.pack(ipadx=60,ipady=4)
    ebutton=Button(eframe,text='send',bg='lightgreen',command=makesit,activebackground='peach puff',relief=GROOVE,font='Times')
    ebutton.pack(ipadx=20)




elabel3=None
elabel2=None
elabel3made=False;elabel2made=False
def sendit():
    global eframe,e1,e2,e3,e4,e5,e6,elabel3made,elabel2made,elabel3,elabel2,var,ethread
    ethread='started'
    efrom=e1.get()
    epassword=e2.get()
    eto=e3.get()
    esubject=e4.get()
    etext=e5.get(index1=1.0,index2=END)
    if var:
       eattach=e6.get()
    if elabel3!=None:
        try:
           elabel3.destroy()
        except:
           elabel3=None
    if elabel2!=None:
       try:
         elabel2.destroy()
       except:
         elabel2=None
    elabel=Label(eframe,text='Sending mail. Please Wait.')
    elabel.pack()
    try:
       if var:
          send_mail(efrom, eto,esubject, etext,files=[eattach],username=efrom, password=epassword)
       else:
          send_mail2(efrom,eto,efrom,epassword,esubject,etext)
       elabel.destroy()
       elabel2=Label(eframe,text='Sent Successfully!',font='Times')
       elabel2.pack()
       elabel2made=True
    except:
       elabel.destroy()
       elabel3=Label(eframe,text='Sending Failed')
       elabel3.pack()
       elabel3made=True



def send_mail2(sender,reciever,username,password,subject,text):
   message="subject: {} \n{}".format(subject,text)
   a=smtplib.SMTP('smtp.gmail.com',25)
   a.ehlo()
   a.starttls()
   a.login(username,password)
   a.sendmail(sender,reciever,message)


def makesit():
   sit=threading.Thread(target=sendit)
   sit.start()


searchlabel=None
thisdict2=None
muButton=None
myButtontwo=None
thinitbut=None
def main():
   global root,frame1,myButton,thisdict2,myButtontwo,thinitbut,frameexists,searchlabel
   root = Tk()
   root.title("LookUp")
   root.configure(background="black")
   myButton = Button(root, text="Search inside Documents",bg="SteelBlue1",relief=FLAT,activebackground="gray53",width=18)
   myButton.pack(side=TOP,pady=1)
   myButtontwo = Button(root, text="Search",bg="SteelBlue1",relief=FLAT,activebackground="gray53",width=14,command=buttonPushed)
   myButtontwo.pack(side=TOP,pady=1)

   createTextBox(root)
   frame1=Frame(root,height=800,width=600)
   frame1.pack(side=TOP)
   frame1.config(background='peach puff')
   frame1.pack_propagate(False)
   searchlabel=Message(frame1,bg='peach puff',fg='blue',text="`I Search Everywhere, I Search Smart. \n `Type an artist's name and their creations will be shown ",width=1000,font='Times 10')
   searchlabel.pack(side=TOP)

   mymenu=Menu(root)
   frameexists=True

   utils=Menu(mymenu,tearoff=0)
   thinitbut=threading.Thread(target=initbut)
   thinitbut.start()
   utils.add_command(label='Update Database',activebackground='SteelBlue1',background='peach puff',command=upcommand)
   utils.add_separator()
   utils.add_command(label='Quick email',activebackground='SteelBlue1',background='peach puff',command=init_mail)
   utils.add_separator()
   utils.add_command(label='Remove Duplicate Files',activebackground='SteelBlue1',background='peach puff',command=init_dup)
   utils.add_separator()
   utils.add_command(label='Text File Punctuation Remover',activebackground='SteelBlue1',background='peach puff',command=inityughere)
   utils.add_separator()
   mymenu.add_cascade(label='Additional Utilities',menu=utils)
   root.config(menu=mymenu)
   myButton.bind('<Button-1>',search2)
   root.mainloop()

main()

from autoupdateDATA import *
